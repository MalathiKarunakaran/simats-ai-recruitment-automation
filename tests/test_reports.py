import io
from datetime import date, datetime, timedelta, timezone

import openpyxl
from pptx import Presentation

from app.models.enums import ApplicationStatusEnum, InterviewTypeEnum, StaffRoleCategoryEnum, UserRoleEnum
from app.services import interviews as interviews_service

from tests.conftest import auth_headers


def _report(client, actor, report_type, **params):
    return client.get(f"/api/v1/reports/{report_type}", headers=auth_headers(client, actor), params=params)


def _export(client, actor, report_type, **params):
    return client.get(f"/api/v1/reports/{report_type}/export", headers=auth_headers(client, actor), params=params)


def _weekly_status(client, actor, **params):
    return client.get("/api/v1/reports/weekly-status", headers=auth_headers(client, actor), params=params)


def _weekly_status_export(client, actor, **params):
    return client.get(
        "/api/v1/reports/weekly-status/export", headers=auth_headers(client, actor), params=params
    )


def test_recruitment_funnel_report_counts(client, published_vacancy_factory, application_factory):
    vacancy = published_vacancy_factory(campus_code="SSE", slot_count=2)
    application_factory(vacancy.job_posting, recorded_by=vacancy.hr_admin)
    application_factory(vacancy.job_posting, recorded_by=vacancy.hr_admin)

    response = _report(client, vacancy.hr_admin, "recruitment-funnel")
    assert response.status_code == 200
    rows = response.json()["rows"]
    row = next(r for r in rows if r["campus_code"] == "SSE" and r["status"] == "APPLIED")
    assert row["role_category"] == "TEACHING"
    assert row["count"] == 2


def test_recruitment_funnel_report_role_category_filter(client, published_vacancy_factory, application_factory):
    teaching_vacancy = published_vacancy_factory(
        campus_code="SSE", slot_count=1, role_category=StaffRoleCategoryEnum.TEACHING
    )
    non_teaching_vacancy = published_vacancy_factory(
        campus_code="SSE", slot_count=1, role_category=StaffRoleCategoryEnum.NON_TEACHING
    )
    application_factory(teaching_vacancy.job_posting, recorded_by=teaching_vacancy.hr_admin)
    application_factory(non_teaching_vacancy.job_posting, recorded_by=non_teaching_vacancy.hr_admin)

    response = _report(client, teaching_vacancy.hr_admin, "recruitment-funnel", role_category="TEACHING")
    assert response.status_code == 200
    rows = response.json()["rows"]
    assert rows
    assert {r["role_category"] for r in rows} == {"TEACHING"}
    assert sum(r["count"] for r in rows if r["campus_code"] == "SSE") == 1


def test_campus_role_hiring_report(client, published_vacancy_factory, hired_employee_factory):
    vacancy = published_vacancy_factory(campus_code="SSE", slot_count=1)
    hired_employee_factory(vacancy)

    response = _report(client, vacancy.hr_admin, "campus-role-hiring")
    rows = response.json()["rows"]
    row = next(r for r in rows if r["campus_code"] == "SSE")
    assert row["role_category"] == "TEACHING"
    assert row["hired_count"] == 1


def test_offer_and_joining_reports(client, published_vacancy_factory, hired_employee_factory):
    vacancy = published_vacancy_factory(campus_code="SSE", slot_count=1)
    hired_employee_factory(vacancy)

    offers_resp = _report(client, vacancy.hr_admin, "offers")
    assert offers_resp.status_code == 200
    offer_row = next(r for r in offers_resp.json()["rows"] if r["campus_code"] == "SSE")
    assert offer_row["role_category"] == "TEACHING"
    assert offer_row["status"] == "ACCEPTED"
    assert offer_row["count"] == 1

    joining_resp = _report(client, vacancy.hr_admin, "joining")
    assert joining_resp.status_code == 200
    joining_row = next(r for r in joining_resp.json()["rows"] if r["campus_code"] == "SSE")
    assert joining_row["role_category"] == "TEACHING"
    assert joining_row["onboarding_status"] == "COMPLETE"
    assert joining_row["count"] == 1


def test_interview_report_counts(client, published_vacancy_factory, application_factory, user_factory, db_session):
    vacancy = published_vacancy_factory(campus_code="SSE", slot_count=1)
    application = application_factory(vacancy.job_posting, recorded_by=vacancy.hr_admin)
    panel_member = user_factory(UserRoleEnum.INTERVIEW_PANEL_MEMBER, campus_code="SSE")

    interviews_service.schedule_interview(
        db_session,
        application=application,
        interview_type=InterviewTypeEnum.TECHNICAL,
        scheduled_at=datetime.now(timezone.utc) + timedelta(days=1),
        duration_minutes=45,
        meeting_link="https://meet.example.com/test",
        location=None,
        notes=None,
        panel_member_ids=[panel_member.id],
        actor=vacancy.hr_admin,
    )

    response = _report(client, vacancy.hr_admin, "interviews")
    assert response.status_code == 200
    row = next(r for r in response.json()["rows"] if r["campus_code"] == "SSE")
    assert row["role_category"] == "TEACHING"
    assert row["status"] == "SCHEDULED"
    assert row["interview_type"] == "TECHNICAL"
    assert row["count"] == 1


def test_vacancy_report_counts(client, published_vacancy_factory):
    vacancy = published_vacancy_factory(campus_code="SSE", slot_count=1)
    response = _report(client, vacancy.hr_admin, "vacancies")
    rows = response.json()["rows"]
    row = next(r for r in rows if r["campus_code"] == "SSE" and r["status"] == "PUBLISHED")
    assert row["count"] == 1


def test_time_to_hire_report_exact_average(client, published_vacancy_factory, hired_employee_factory):
    vacancy = published_vacancy_factory(campus_code="SSE", slot_count=1)
    hired_employee_factory(
        vacancy, applied_at=datetime(2026, 1, 1, tzinfo=timezone.utc), joining_date=date(2026, 1, 11)
    )

    response = _report(client, vacancy.hr_admin, "time-to-hire")
    rows = response.json()["rows"]
    row = next(r for r in rows if r["campus_code"] == "SSE")
    assert row["avg_days"] == 10.0
    assert row["hired_count"] == 1


def test_report_respects_campus_scoping(client, published_vacancy_factory, application_factory):
    vacancy_sse = published_vacancy_factory(campus_code="SSE", slot_count=1)
    vacancy_scad = published_vacancy_factory(campus_code="SCAD", slot_count=1)
    application_factory(vacancy_sse.job_posting, recorded_by=vacancy_sse.hr_admin)
    application_factory(vacancy_scad.job_posting, recorded_by=vacancy_scad.hr_admin)

    hod_response = _report(client, vacancy_sse.hod, "recruitment-funnel")
    hod_codes = {row["campus_code"] for row in hod_response.json()["rows"]}
    assert hod_codes == {"SSE"}

    global_response = _report(client, vacancy_sse.hr_admin, "recruitment-funnel")
    global_codes = {row["campus_code"] for row in global_response.json()["rows"]}
    assert {"SSE", "SCAD"} <= global_codes


def test_unknown_report_type_returns_404(client, published_vacancy_factory):
    vacancy = published_vacancy_factory(campus_code="SSE", slot_count=1)
    response = _report(client, vacancy.hr_admin, "not-a-real-report")
    assert response.status_code == 404


def test_candidate_role_forbidden(client, user_factory):
    candidate_user = user_factory(UserRoleEnum.CANDIDATE)
    response = _report(client, candidate_user, "vacancies")
    assert response.status_code == 403


def test_export_xlsx_is_valid_and_matches_headers(client, published_vacancy_factory):
    vacancy = published_vacancy_factory(campus_code="SSE", slot_count=1)
    response = _export(client, vacancy.hr_admin, "vacancies")
    assert response.status_code == 200
    assert response.headers["content-type"] == (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    )

    workbook = openpyxl.load_workbook(io.BytesIO(response.content))
    sheet = workbook.active
    header_row = [cell.value for cell in sheet[4]]
    assert header_row == ["campus_code", "role_category", "status", "count"]


def test_export_xlsx_headers_include_role_category_for_newly_grouped_reports(
    client, published_vacancy_factory, application_factory
):
    """recruitment-funnel/interviews/offers/joining didn't group by
    role_category before this phase -- confirm their xlsx export header rows
    now include it too, alongside the export data cell in each row."""
    vacancy = published_vacancy_factory(campus_code="SSE", slot_count=1)
    application_factory(vacancy.job_posting, recorded_by=vacancy.hr_admin)

    funnel_response = _export(client, vacancy.hr_admin, "recruitment-funnel")
    assert funnel_response.status_code == 200
    funnel_workbook = openpyxl.load_workbook(io.BytesIO(funnel_response.content))
    funnel_header = [cell.value for cell in funnel_workbook.active[4]]
    assert funnel_header == ["campus_code", "role_category", "status", "count"]
    funnel_data_row = [cell.value for cell in funnel_workbook.active[5]]
    assert "TEACHING" in funnel_data_row

    interviews_response = _export(client, vacancy.hr_admin, "interviews")
    assert interviews_response.status_code == 200
    interviews_workbook = openpyxl.load_workbook(io.BytesIO(interviews_response.content))
    interviews_header = [cell.value for cell in interviews_workbook.active[4]]
    assert interviews_header == ["campus_code", "role_category", "status", "interview_type", "count"]

    offers_response = _export(client, vacancy.hr_admin, "offers")
    assert offers_response.status_code == 200
    offers_workbook = openpyxl.load_workbook(io.BytesIO(offers_response.content))
    offers_header = [cell.value for cell in offers_workbook.active[4]]
    assert offers_header == ["campus_code", "role_category", "status", "count"]

    joining_response = _export(client, vacancy.hr_admin, "joining")
    assert joining_response.status_code == 200
    joining_workbook = openpyxl.load_workbook(io.BytesIO(joining_response.content))
    joining_header = [cell.value for cell in joining_workbook.active[4]]
    assert joining_header == ["campus_code", "role_category", "onboarding_status", "count"]


def test_ad_briefing_json_and_export(client, published_vacancy_factory, hired_employee_factory):
    vacancy = published_vacancy_factory(campus_code="SSE", slot_count=2)
    hired_employee_factory(vacancy)

    json_response = client.get(
        "/api/v1/reports/ad-briefing", headers=auth_headers(client, vacancy.hr_admin)
    )
    assert json_response.status_code == 200
    body = json_response.json()
    assert "kpi_headline" in body
    breakdown_row = next(r for r in body["campus_role_breakdown"] if r["campus_code"] == "SSE")
    assert breakdown_row["hired"] == 1

    export_response = client.get(
        "/api/v1/reports/ad-briefing/export", headers=auth_headers(client, vacancy.hr_admin)
    )
    assert export_response.status_code == 200
    assert export_response.headers["content-type"] == (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

    presentation = Presentation(io.BytesIO(export_response.content))
    assert len(presentation.slides) == 1
    all_text = " ".join(
        shape.text_frame.text
        for shape in presentation.slides[0].shapes
        if shape.has_text_frame
    )
    assert "AD Briefing" in all_text


def test_ad_briefing_defaults_to_today_period_label(client, published_vacancy_factory, hired_employee_factory):
    vacancy = published_vacancy_factory(campus_code="SSE", slot_count=1)
    hired_employee_factory(vacancy)

    response = client.get("/api/v1/reports/ad-briefing", headers=auth_headers(client, vacancy.hr_admin))
    assert response.status_code == 200
    assert response.json()["period_label"] == "Today"


def test_ad_briefing_date_range_narrows_kpi_headline_and_sets_period_label(
    client, published_vacancy_factory, application_factory
):
    vacancy = published_vacancy_factory(campus_code="SSE", slot_count=2)
    application_factory(vacancy.job_posting, recorded_by=vacancy.hr_admin)

    today = date.today()
    week_ago = today - timedelta(days=7)

    response = client.get(
        "/api/v1/reports/ad-briefing",
        headers=auth_headers(client, vacancy.hr_admin),
        params={"start_date": week_ago.isoformat(), "end_date": today.isoformat()},
    )
    assert response.status_code == 200
    body = response.json()
    assert body["kpi_headline"]["total_applications"] == 1
    assert body["period_label"] == f"{week_ago.isoformat()} to {today.isoformat()}"

    future_start = today + timedelta(days=1)
    excluding_response = client.get(
        "/api/v1/reports/ad-briefing",
        headers=auth_headers(client, vacancy.hr_admin),
        params={"start_date": future_start.isoformat(), "end_date": future_start.isoformat()},
    )
    assert excluding_response.status_code == 200
    assert excluding_response.json()["kpi_headline"]["total_applications"] == 0


def test_ad_briefing_rejects_start_after_end_date(client, published_vacancy_factory):
    vacancy = published_vacancy_factory(campus_code="SSE", slot_count=1)
    today = date.today()

    response = client.get(
        "/api/v1/reports/ad-briefing",
        headers=auth_headers(client, vacancy.hr_admin),
        params={"start_date": today.isoformat(), "end_date": (today - timedelta(days=1)).isoformat()},
    )
    assert response.status_code == 422


def test_recruitment_funnel_report_date_range_narrows_by_applied_at(
    client, published_vacancy_factory, hired_employee_factory
):
    vacancy = published_vacancy_factory(campus_code="SSE", slot_count=2)
    hired_employee_factory(
        vacancy, applied_at=datetime(2026, 1, 1, tzinfo=timezone.utc), joining_date=date(2026, 1, 11)
    )

    in_range = _report(
        client, vacancy.hr_admin, "recruitment-funnel", start_date="2025-12-25", end_date="2026-01-05"
    )
    assert in_range.status_code == 200
    assert sum(r["count"] for r in in_range.json()["rows"] if r["campus_code"] == "SSE") == 1

    out_of_range = _report(
        client, vacancy.hr_admin, "recruitment-funnel", start_date="2026-02-01", end_date="2026-02-28"
    )
    assert out_of_range.status_code == 200
    assert sum(r["count"] for r in out_of_range.json()["rows"] if r["campus_code"] == "SSE") == 0


def test_campus_role_hiring_report_date_range_narrows_by_date_of_joining(
    client, published_vacancy_factory, hired_employee_factory
):
    vacancy = published_vacancy_factory(campus_code="SSE", slot_count=1)
    hired_employee_factory(vacancy, joining_date=date(2026, 3, 15))

    in_range = _report(
        client, vacancy.hr_admin, "campus-role-hiring", start_date="2026-03-01", end_date="2026-03-31"
    )
    row = next((r for r in in_range.json()["rows"] if r["campus_code"] == "SSE"), None)
    assert row is not None and row["hired_count"] == 1

    out_of_range = _report(
        client, vacancy.hr_admin, "campus-role-hiring", start_date="2026-04-01", end_date="2026-04-30"
    )
    assert not any(r["campus_code"] == "SSE" for r in out_of_range.json()["rows"])


def test_interview_report_date_range_narrows_by_scheduled_at(
    client, published_vacancy_factory, application_factory, user_factory, db_session
):
    vacancy = published_vacancy_factory(campus_code="SSE", slot_count=1)
    application = application_factory(vacancy.job_posting, recorded_by=vacancy.hr_admin)
    panel_member = user_factory(UserRoleEnum.INTERVIEW_PANEL_MEMBER, campus_code="SSE")

    scheduled_at = datetime(2026, 5, 10, 10, 0, tzinfo=timezone.utc)
    interviews_service.schedule_interview(
        db_session,
        application=application,
        interview_type=InterviewTypeEnum.TECHNICAL,
        scheduled_at=scheduled_at,
        duration_minutes=45,
        meeting_link="https://meet.example.com/test",
        location=None,
        notes=None,
        panel_member_ids=[panel_member.id],
        actor=vacancy.hr_admin,
    )

    in_range = _report(client, vacancy.hr_admin, "interviews", start_date="2026-05-01", end_date="2026-05-31")
    assert sum(r["count"] for r in in_range.json()["rows"] if r["campus_code"] == "SSE") == 1

    out_of_range = _report(client, vacancy.hr_admin, "interviews", start_date="2026-06-01", end_date="2026-06-30")
    assert sum(r["count"] for r in out_of_range.json()["rows"] if r["campus_code"] == "SSE") == 0


def test_offer_report_date_range_narrows_by_created_at(client, published_vacancy_factory, hired_employee_factory):
    vacancy = published_vacancy_factory(campus_code="SSE", slot_count=1)
    hired_employee_factory(vacancy)

    today = date.today()
    in_range = _report(
        client, vacancy.hr_admin, "offers", start_date=(today - timedelta(days=1)).isoformat(), end_date=today.isoformat()
    )
    assert sum(r["count"] for r in in_range.json()["rows"] if r["campus_code"] == "SSE") == 1

    past_start = today - timedelta(days=30)
    past_end = today - timedelta(days=20)
    out_of_range = _report(
        client, vacancy.hr_admin, "offers", start_date=past_start.isoformat(), end_date=past_end.isoformat()
    )
    assert sum(r["count"] for r in out_of_range.json()["rows"] if r["campus_code"] == "SSE") == 0


def test_joining_report_date_range_narrows_by_created_at_not_actual_joining_date(
    client, published_vacancy_factory, hired_employee_factory
):
    vacancy = published_vacancy_factory(campus_code="SSE", slot_count=1)
    hired_employee_factory(vacancy)

    today = date.today()
    in_range = _report(
        client, vacancy.hr_admin, "joining", start_date=(today - timedelta(days=1)).isoformat(), end_date=today.isoformat()
    )
    row = next(r for r in in_range.json()["rows"] if r["campus_code"] == "SSE")
    assert row["onboarding_status"] == "COMPLETE"
    assert row["count"] == 1

    past_start = today - timedelta(days=30)
    past_end = today - timedelta(days=20)
    out_of_range = _report(
        client, vacancy.hr_admin, "joining", start_date=past_start.isoformat(), end_date=past_end.isoformat()
    )
    assert not any(r["campus_code"] == "SSE" for r in out_of_range.json()["rows"])


def test_vacancy_report_date_range_narrows_by_created_at(client, published_vacancy_factory):
    vacancy = published_vacancy_factory(campus_code="SSE", slot_count=1)

    today = date.today()
    in_range = _report(
        client, vacancy.hr_admin, "vacancies", start_date=(today - timedelta(days=1)).isoformat(), end_date=today.isoformat()
    )
    row = next(r for r in in_range.json()["rows"] if r["campus_code"] == "SSE" and r["status"] == "PUBLISHED")
    assert row["count"] == 1

    past_start = today - timedelta(days=30)
    past_end = today - timedelta(days=20)
    out_of_range = _report(
        client, vacancy.hr_admin, "vacancies", start_date=past_start.isoformat(), end_date=past_end.isoformat()
    )
    assert not any(
        r["campus_code"] == "SSE" and r["status"] == "PUBLISHED" for r in out_of_range.json()["rows"]
    )


def test_time_to_hire_report_date_range_narrows_by_actual_joining_date(
    client, published_vacancy_factory, hired_employee_factory
):
    vacancy = published_vacancy_factory(campus_code="SSE", slot_count=1)
    hired_employee_factory(
        vacancy, applied_at=datetime(2026, 1, 1, tzinfo=timezone.utc), joining_date=date(2026, 1, 11)
    )

    in_range = _report(client, vacancy.hr_admin, "time-to-hire", start_date="2026-01-01", end_date="2026-01-31")
    row = next(r for r in in_range.json()["rows"] if r["campus_code"] == "SSE")
    assert row["avg_days"] == 10.0
    assert row["hired_count"] == 1

    out_of_range = _report(client, vacancy.hr_admin, "time-to-hire", start_date="2026-02-01", end_date="2026-02-28")
    assert not any(r["campus_code"] == "SSE" for r in out_of_range.json()["rows"])


def test_report_rejects_start_after_end_date(client, published_vacancy_factory):
    vacancy = published_vacancy_factory(campus_code="SSE", slot_count=1)
    today = date.today()

    response = _report(
        client,
        vacancy.hr_admin,
        "vacancies",
        start_date=today.isoformat(),
        end_date=(today - timedelta(days=1)).isoformat(),
    )
    assert response.status_code == 422


def test_ad_briefing_export_reflects_period_label_in_slide_text(client, published_vacancy_factory):
    vacancy = published_vacancy_factory(campus_code="SSE", slot_count=1)
    today = date.today()

    response = client.get(
        "/api/v1/reports/ad-briefing/export",
        headers=auth_headers(client, vacancy.hr_admin),
        params={"start_date": today.isoformat(), "end_date": today.isoformat()},
    )
    assert response.status_code == 200
    presentation = Presentation(io.BytesIO(response.content))
    all_text = " ".join(
        shape.text_frame.text for shape in presentation.slides[0].shapes if shape.has_text_frame
    )
    assert f"Interviews ({today.isoformat()})" in all_text


def test_weekly_status_teaching_row_selected(
    client, published_vacancy_factory, application_factory, db_session
):
    vacancy = published_vacancy_factory(
        campus_code="SSE", slot_count=1, role_category=StaffRoleCategoryEnum.TEACHING
    )
    application = application_factory(vacancy.job_posting, recorded_by=vacancy.hr_admin)
    today = date.today()
    application.panel_result = "Selected"
    application.interview_scheduled_date = today
    db_session.flush()

    response = _weekly_status(
        client, vacancy.hr_admin, start_date=today.isoformat(), end_date=today.isoformat()
    )
    assert response.status_code == 200
    body = response.json()
    row = next(r for r in body["teaching_rows"] if r["group_label"] == "SSE")
    assert row["attended"] == 1
    assert row["selected"] == 1
    assert row["waiting"] == 0
    assert row["rejected"] == 0
    assert row["upcoming_join"] is None
    assert row["joined"] is None
    assert body["total_interviewed"] == 1
    assert body["total_selected"] == 1
    assert body["selection_rate_pct"] == 100.0 or body["selection_rate_pct"] == 100


def test_weekly_status_non_teaching_upcoming_join(
    client, published_vacancy_factory, application_factory, db_session
):
    vacancy = published_vacancy_factory(
        campus_code="SSE", slot_count=1, role_category=StaffRoleCategoryEnum.NON_TEACHING
    )
    application = application_factory(vacancy.job_posting, recorded_by=vacancy.hr_admin)
    today = date.today()
    application.expected_joining_date = today + timedelta(days=7)
    application.status = ApplicationStatusEnum.OFFER_ACCEPTED
    db_session.flush()

    response = _weekly_status(
        client, vacancy.hr_admin, start_date=today.isoformat(), end_date=today.isoformat()
    )
    assert response.status_code == 200
    body = response.json()
    row = next(r for r in body["non_teaching_rows"] if r["group_label"] == vacancy.vacancy_request.position_title)
    assert row["upcoming_join"] == 1
    assert row["joined"] == 0


def test_weekly_status_joined_this_week(
    client, published_vacancy_factory, application_factory, db_session
):
    vacancy = published_vacancy_factory(
        campus_code="SSE", slot_count=1, role_category=StaffRoleCategoryEnum.NON_TEACHING
    )
    application = application_factory(vacancy.job_posting, recorded_by=vacancy.hr_admin)
    today = date.today()
    application.actual_joining_date = today
    application.status = ApplicationStatusEnum.JOINED
    db_session.flush()

    response = _weekly_status(
        client, vacancy.hr_admin, start_date=today.isoformat(), end_date=today.isoformat()
    )
    assert response.status_code == 200
    body = response.json()
    assert body["total_joined"] == 1
    assert body["joined_by_category"]["NON_TEACHING"] == 1
    row = next(r for r in body["non_teaching_rows"] if r["group_label"] == vacancy.vacancy_request.position_title)
    assert row["joined"] == 1


def test_weekly_status_requires_start_and_end_date(client, published_vacancy_factory):
    vacancy = published_vacancy_factory(campus_code="SSE", slot_count=1)
    response = _weekly_status(client, vacancy.hr_admin)
    assert response.status_code == 422


def test_weekly_status_rejects_start_after_end_date(client, published_vacancy_factory):
    vacancy = published_vacancy_factory(campus_code="SSE", slot_count=1)
    today = date.today()
    response = _weekly_status(
        client, vacancy.hr_admin, start_date=today.isoformat(), end_date=(today - timedelta(days=1)).isoformat()
    )
    assert response.status_code == 422


def test_weekly_status_export_returns_two_slide_pptx(
    client, published_vacancy_factory, application_factory, db_session
):
    vacancy = published_vacancy_factory(campus_code="SSE", slot_count=1)
    application = application_factory(vacancy.job_posting, recorded_by=vacancy.hr_admin)
    today = date.today()
    application.panel_result = "Selected"
    application.interview_scheduled_date = today
    db_session.flush()

    response = _weekly_status_export(
        client, vacancy.hr_admin, start_date=today.isoformat(), end_date=today.isoformat()
    )
    assert response.status_code == 200
    assert response.headers["content-type"] == (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    assert len(response.content) > 0

    presentation = Presentation(io.BytesIO(response.content))
    assert len(presentation.slides) == 2
    all_text = " ".join(
        shape.text_frame.text
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "Weekly Recruitment Status" in all_text
    assert "Prepared by the Recruitment Office" in all_text


def test_weekly_status_respects_campus_scoping(
    client, published_vacancy_factory, application_factory, db_session
):
    vacancy_sse = published_vacancy_factory(campus_code="SSE", slot_count=1)
    vacancy_scad = published_vacancy_factory(campus_code="SCAD", slot_count=1)
    today = date.today()

    app_sse = application_factory(vacancy_sse.job_posting, recorded_by=vacancy_sse.hr_admin)
    app_sse.panel_result = "Selected"
    app_sse.interview_scheduled_date = today

    app_scad = application_factory(vacancy_scad.job_posting, recorded_by=vacancy_scad.hr_admin)
    app_scad.panel_result = "Selected"
    app_scad.interview_scheduled_date = today
    db_session.flush()

    hod_response = _weekly_status(
        client, vacancy_sse.hod, start_date=today.isoformat(), end_date=today.isoformat()
    )
    assert hod_response.status_code == 200
    hod_labels = {r["group_label"] for r in hod_response.json()["teaching_rows"]}
    assert hod_labels == {"SSE"}

    global_response = _weekly_status(
        client, vacancy_sse.hr_admin, start_date=today.isoformat(), end_date=today.isoformat()
    )
    assert global_response.status_code == 200
    global_labels = {r["group_label"] for r in global_response.json()["teaching_rows"]}
    assert {"SSE", "SCAD"} <= global_labels
