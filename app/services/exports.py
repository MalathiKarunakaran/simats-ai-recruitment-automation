"""File-generation helpers for Module 12's Excel/PPT exports.

Pure functions: given already-computed report/summary data, return raw
bytes. No DB access here -- app/services/reporting.py owns all querying.

PPT branding: the actual current SIMATS AD-meeting PPT template, logo, and
exact brand hex codes weren't available to reconstruct exactly, so
_NAVY/_GOLD below are documented placeholders reproducing the described
structure (navy/gold, single slide, campus x role-category breakdown) --
swap in the real brand assets when available.
"""

import io
from datetime import datetime

from openpyxl import Workbook
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

# Placeholder brand colors -- TODO: replace with the real SIMATS navy/gold
# hex codes and logo once the actual AD-meeting PPT template is available.
_NAVY = RGBColor(0x0B, 0x1F, 0x3A)
_GOLD = RGBColor(0xC9, 0xA2, 0x27)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

_REPORT_FIELDS: dict[str, list[str]] = {
    "recruitment-funnel": ["campus_code", "role_category", "status", "count"],
    "campus-role-hiring": ["campus_code", "role_category", "hired_count"],
    "interviews": ["campus_code", "role_category", "status", "interview_type", "count"],
    "offers": ["campus_code", "role_category", "status", "count"],
    "joining": ["campus_code", "role_category", "onboarding_status", "count"],
    "vacancies": ["campus_code", "role_category", "status", "count"],
    "time-to-hire": ["campus_code", "role_category", "avg_days", "hired_count"],
    "sanctioned-strength-reconciliation": [
        "campus_code",
        "department_name",
        "designation_name",
        "category",
        "approved_strength",
        "working_count",
        "already_requested",
        "over_by",
    ],
    "resignations": ["campus_code", "department_name", "role_category", "count"],
}


def build_report_excel(report_type: str, rows: list[dict], generated_at: datetime, scope_note: str) -> bytes:
    fields = _REPORT_FIELDS[report_type]

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = report_type[:31]

    sheet.append([f"Generated: {generated_at.isoformat()}"])
    sheet.append([f"Scope: {scope_note}"])
    sheet.append([])
    sheet.append(fields)
    for row in rows:
        sheet.append([row.get(field) for field in fields])

    buf = io.BytesIO()
    workbook.save(buf)
    return buf.getvalue()


_DEPARTMENT_EXPORT_HEADERS = [
    "Campus Code",
    "Department Code",
    "Department Name",
    "Category",
    "Parent Group",
    "Description",
    "Active",
]


def build_department_export_excel(rows: list[dict], generated_at: datetime, scope_note: str) -> bytes:
    """Department Master's own dedicated export builder (Department Master
    hardening epic, 2026-08-25) -- not `build_report_excel` above (that one's
    `_REPORT_FIELDS` is keyed to Module 12's own fixed `report_type` strings,
    which Department Master export isn't one of), but the same visual/header
    convention: a `Generated:`/`Scope:` line, a blank row, then a real header
    row, matching every other export this app produces.
    """
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Departments"

    sheet.append([f"Generated: {generated_at.isoformat()}"])
    sheet.append([f"Scope: {scope_note}"])
    sheet.append([])
    sheet.append(_DEPARTMENT_EXPORT_HEADERS)
    for row in rows:
        sheet.append(
            [
                row.get("campus_code"),
                row.get("code"),
                row.get("name"),
                row.get("category"),
                row.get("parent_group"),
                row.get("description"),
                "TRUE" if row.get("is_active") else "FALSE",
            ]
        )

    buf = io.BytesIO()
    workbook.save(buf)
    return buf.getvalue()


_ELIGIBILITY_RULE_EXPORT_HEADERS = [
    "Campus Code",
    "Department Code",
    "Staff Category",
    "Position Title",
    "Regulatory Authority",
    "School/College",
    "Programme/Discipline",
    "Required Qualification Keyword",
    "Minimum Qualification",
    "Minimum Percentage",
    "Required Experience",
    "Required Credential",
    "NET/SET/SLET Required",
    "PhD Required",
    "Professional Registration",
    "Industry Experience",
    "Priority",
    "Effective From",
    "Effective To",
    "Source Regulation",
    "Status",
    "Verification Required",
    "Active",
    "Notes",
]


def build_eligibility_rule_export_excel(rows: list[dict], generated_at: datetime, scope_note: str) -> bytes:
    """EligibilityRule's own dedicated export builder (starter regulatory-
    eligibility-rules feature, backend Phase 1) -- same visual/header
    convention as `build_department_export_excel` above (`Generated:`/
    `Scope:` line, a blank row, then a real header row). Only exports the
    fields most useful for an admin reviewing a regulatory mapping in bulk;
    every column round-trips through the bulk-upload template as well (see
    app/services/eligibility_rule_import.py), so nothing here is export-only.
    """
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Eligibility Rules"

    sheet.append([f"Generated: {generated_at.isoformat()}"])
    sheet.append([f"Scope: {scope_note}"])
    sheet.append([])
    sheet.append(_ELIGIBILITY_RULE_EXPORT_HEADERS)
    for row in rows:
        sheet.append(
            [
                row.get("campus_code"),
                row.get("department_code"),
                row.get("staff_category"),
                row.get("position_title"),
                row.get("regulatory_authority"),
                row.get("school_or_college"),
                row.get("programme_discipline"),
                row.get("required_qualification_keyword"),
                row.get("minimum_qualification"),
                row.get("minimum_percentage"),
                row.get("required_experience"),
                row.get("required_credential"),
                _bool_or_blank(row.get("net_set_required")),
                _bool_or_blank(row.get("phd_required")),
                row.get("professional_registration"),
                row.get("industry_experience"),
                row.get("priority"),
                row.get("effective_from").isoformat() if row.get("effective_from") else None,
                row.get("effective_to").isoformat() if row.get("effective_to") else None,
                row.get("source_regulation"),
                row.get("status"),
                "TRUE" if row.get("verification_required") else "FALSE",
                "TRUE" if row.get("is_active") else "FALSE",
                row.get("notes"),
            ]
        )

    buf = io.BytesIO()
    workbook.save(buf)
    return buf.getvalue()


def _bool_or_blank(value: bool | None) -> str | None:
    if value is None:
        return None
    return "TRUE" if value else "FALSE"


def build_ad_briefing_pptx(summary: dict) -> bytes:
    presentation = Presentation()
    presentation.slide_width = Inches(13.33)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # blank layout

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, presentation.slide_width, presentation.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = _NAVY
    background.line.fill.background()
    background.shadow.inherit = False

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
    title_frame = title_box.text_frame
    title_frame.text = "SIMATS Recruitment — AD Briefing"
    title_run = title_frame.paragraphs[0].runs[0]
    title_run.font.size = Pt(32)
    title_run.font.bold = True
    title_run.font.color.rgb = _GOLD

    generated_at = summary["generated_at"]
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.4))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = f"{summary['scope_note']} — generated {generated_at.isoformat()}"
    subtitle_run = subtitle_frame.paragraphs[0].runs[0]
    subtitle_run.font.size = Pt(14)
    subtitle_run.font.color.rgb = _WHITE

    kpi = summary["kpi_headline"]
    period_label = summary["period_label"]
    # vacancy_closure_rate_pct is None (not 0.0) when there's nothing
    # APPROVED-or-beyond to compute it from -- render that as prose, not a
    # literal "None%" (CLAUDE.md B1).
    closure_rate_text = (
        f"{kpi['vacancy_closure_rate_pct']}%" if kpi["vacancy_closure_rate_pct"] is not None else "Not enough data yet"
    )
    kpi_text = (
        f"Applications: {kpi['total_applications']}  |  Open positions: {kpi['open_positions']}  |  "
        f"Interviews ({period_label}): {kpi['interviews_today']}  |  Joinings ({period_label}): {kpi['joinings_today']}  |  "
        f"Offers pending: {kpi['offers_pending']}  |  Closure rate: {closure_rate_text}"
    )
    kpi_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(12.3), Inches(0.6))
    kpi_frame = kpi_box.text_frame
    kpi_frame.word_wrap = True
    kpi_frame.text = kpi_text
    kpi_run = kpi_frame.paragraphs[0].runs[0]
    kpi_run.font.size = Pt(14)
    kpi_run.font.color.rgb = _GOLD

    breakdown = summary["campus_role_breakdown"]
    table_rows = max(len(breakdown) + 1, 2)
    table_shape = slide.shapes.add_table(
        table_rows, 5, Inches(0.5), Inches(2.5), Inches(12.3), Inches(4.5)
    )
    table = table_shape.table
    headers = ["Campus", "Role Category", "Open", "In Pipeline", "Hired"]
    for col_index, header in enumerate(headers):
        cell = table.cell(0, col_index)
        cell.text = header
        cell.text_frame.paragraphs[0].runs[0].font.bold = True

    for row_index, entry in enumerate(breakdown, start=1):
        values = [
            entry["campus_code"],
            entry["role_category"],
            str(entry["open_positions"]),
            str(entry["in_pipeline"]),
            str(entry["hired"]),
        ]
        for col_index, value in enumerate(values):
            table.cell(row_index, col_index).text = value

    buf = io.BytesIO()
    presentation.save(buf)
    return buf.getvalue()


def _add_blank_slide(presentation: Presentation):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, presentation.slide_width, presentation.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = _NAVY
    background.line.fill.background()
    background.shadow.inherit = False
    return slide


def _add_textbox(slide, left, top, width, height, text: str, size: int, color: RGBColor, bold: bool = False):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.text = text
    run = frame.paragraphs[0].runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _add_status_table(slide, headers: list[str], rows: list[list[str]], left, top, width, height):
    table_shape = slide.shapes.add_table(max(len(rows) + 1, 2), len(headers), left, top, width, height)
    table = table_shape.table
    for col_index, header in enumerate(headers):
        cell = table.cell(0, col_index)
        cell.text = header
        cell.text_frame.paragraphs[0].runs[0].font.bold = True
    for row_index, values in enumerate(rows, start=1):
        for col_index, value in enumerate(values):
            cell = table.cell(row_index, col_index)
            cell.text = value
            if row_index == len(rows):
                # Rows list always ends with the caller-supplied Total row.
                cell.text_frame.paragraphs[0].runs[0].font.bold = True
    return table_shape


def build_weekly_status_pptx(data: dict) -> bytes:
    presentation = Presentation()
    presentation.slide_width = Inches(13.33)
    presentation.slide_height = Inches(7.5)

    # Slide 1: title/subtitle, top-line KPIs, Teaching Staff table.
    slide1 = _add_blank_slide(presentation)
    _add_textbox(
        slide1, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9),
        "Weekly Recruitment Status", 32, _GOLD, bold=True,
    )
    _add_textbox(
        slide1, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.4),
        f"Teaching & Non-Teaching Staff | SIMATS Engineering — {data['period_label']}", 14, _WHITE,
    )

    kpi_text = (
        f"Total Interviewed: {data['total_interviewed']}  |  Selected: {data['total_selected']}  |  "
        f"Waiting List: {data['total_waiting']}  |  Rejected: {data['total_rejected']}  |  "
        f"Joined This Week: {data['total_joined']}"
    )
    _add_textbox(slide1, Inches(0.5), Inches(1.7), Inches(12.3), Inches(0.6), kpi_text, 14, _GOLD)

    _add_textbox(
        slide1, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.4),
        "Teaching Staff (TS) — Campus-wise", 16, _WHITE, bold=True,
    )
    teaching_rows = data["teaching_rows"]
    teaching_totals = {
        "attended": sum(r["attended"] for r in teaching_rows),
        "selected": sum(r["selected"] for r in teaching_rows),
        "waiting": sum(r["waiting"] for r in teaching_rows),
        "rejected": sum(r["rejected"] for r in teaching_rows),
    }
    teaching_table_rows = [
        [r["group_label"], str(r["attended"]), str(r["selected"]), str(r["waiting"]), str(r["rejected"])]
        for r in teaching_rows
    ] + [
        [
            "Total",
            str(teaching_totals["attended"]),
            str(teaching_totals["selected"]),
            str(teaching_totals["waiting"]),
            str(teaching_totals["rejected"]),
        ]
    ]
    _add_status_table(
        slide1,
        ["College", "Attended", "Selected", "Waiting", "Rejected"],
        teaching_table_rows,
        Inches(0.5), Inches(2.8), Inches(12.3), Inches(4.0),
    )

    # Slide 2: Joined This Week category tiles, Non-Teaching Staff table, footer.
    slide2 = _add_blank_slide(presentation)
    _add_textbox(
        slide2, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.4),
        "Joined This Week — Category-wise", 20, _GOLD, bold=True,
    )
    joined_by_category = data["joined_by_category"]
    category_text = (
        f"Teaching: {joined_by_category.get('TEACHING', 0)}  |  "
        f"Non-Teaching: {joined_by_category.get('NON_TEACHING', 0)}  |  "
        f"Housekeeping: {joined_by_category.get('HOUSEKEEPING', 0)}"
    )
    _add_textbox(slide2, Inches(0.5), Inches(0.9), Inches(12.3), Inches(0.5), category_text, 16, _WHITE)

    _add_textbox(
        slide2, Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.4),
        "Non-Teaching Staff (NTS) — Department-wise", 16, _WHITE, bold=True,
    )
    non_teaching_rows = data["non_teaching_rows"]
    nts_totals = {
        "attended": sum(r["attended"] for r in non_teaching_rows),
        "selected": sum(r["selected"] for r in non_teaching_rows),
        "waiting": sum(r["waiting"] for r in non_teaching_rows),
        "rejected": sum(r["rejected"] for r in non_teaching_rows),
        "upcoming_join": sum(r["upcoming_join"] or 0 for r in non_teaching_rows),
        "joined": sum(r["joined"] or 0 for r in non_teaching_rows),
    }
    nts_table_rows = [
        [
            r["group_label"],
            str(r["attended"]),
            str(r["selected"]),
            str(r["waiting"]),
            str(r["rejected"]),
            str(r["upcoming_join"] or 0),
            str(r["joined"] or 0),
        ]
        for r in non_teaching_rows
    ] + [
        [
            "Total",
            str(nts_totals["attended"]),
            str(nts_totals["selected"]),
            str(nts_totals["waiting"]),
            str(nts_totals["rejected"]),
            str(nts_totals["upcoming_join"]),
            str(nts_totals["joined"]),
        ]
    ]
    _add_status_table(
        slide2,
        ["Department", "Attended", "Selected", "Waiting", "Rejected", "Upcoming Join", "Joined"],
        nts_table_rows,
        Inches(0.5), Inches(2.1), Inches(12.3), Inches(4.2),
    )

    selection_rate_pct = data["selection_rate_pct"]
    rate_text = "N/A" if selection_rate_pct is None else f"{selection_rate_pct}%"
    _add_textbox(
        slide2, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5),
        f"Overall selection rate: {rate_text} | Prepared by the Recruitment Office", 14, _GOLD,
    )

    buf = io.BytesIO()
    presentation.save(buf)
    return buf.getvalue()
